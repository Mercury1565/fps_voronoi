"""Build a simple, near-monochrome discussion deck for FPS-Voronoi.

Focus: today's work — real-LiDAR integration, the frame-to-frame temporal loop,
and the Chamfer quality metric. Theory-oriented; no code stubs.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))

# ── palette: grayscale only ──────────────────────────────────────────────────
INK = RGBColor(0x1A, 0x1A, 0x1A)   # near-black body
MUTE = RGBColor(0x70, 0x70, 0x70)  # muted gray
RULE = RGBColor(0xBF, 0xBF, 0xBF)  # thin divider
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xF2, 0xF2, 0xF2)  # mono code-panel fill

FONT = "Calibri"
MONO = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)  # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def _set(tf_para, text, size, *, bold=False, color=INK, font=FONT, align=None):
    tf_para.text = text if text else " "
    r = tf_para.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    if align is not None:
        tf_para.alignment = align


def _rule(slide, x, y, w):
    """Thin horizontal divider line."""
    from pptx.enum.shapes import MSO_SHAPE
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(1.2))
    ln.fill.solid()
    ln.fill.fore_color.rgb = RULE
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def header(slide, title, kicker=None):
    """Standard slide header: small kicker + title + divider."""
    left, top, width = Inches(0.7), Inches(0.45), Inches(12.0)
    if kicker:
        kb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
        _set(kb.text_frame.paragraphs[0], kicker.upper(), 12, bold=True, color=MUTE)
        top = Inches(0.78)
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    _set(tb.text_frame.paragraphs[0], title, 28, bold=True, color=INK)
    _rule(slide, left, Inches(1.55), Inches(12.0))


def bullets(slide, items, *, left=Inches(0.7), top=Inches(1.8),
            width=Inches(12.0), height=Inches(5.2), size=18, gap=8):
    """items: list of (level, text) or (level, text, 'mute')."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lvl, text = item[0], item[1]
        muted = len(item) > 2 and item[2] == "mute"
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(gap)
        bullet = ("•  " if lvl == 0 else "–  ") if text else ""
        _set(p, bullet + text, size - (2 if lvl else 0),
             color=MUTE if muted else INK)
    return tb


def fit_image(slide, path, *, max_w, max_h, top, center_x=True, left=None):
    """Place an image scaled to fit (max_w, max_h), preserving aspect ratio."""
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = Emu(int(iw * scale)), Emu(int(ih * scale))
    if center_x:
        left = int((SW - w) / 2)
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    return w, h


def caption(slide, text, top, *, width=Inches(12.0), left=Inches(0.7)):
    cb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    _set(cb.text_frame.paragraphs[0], text, 12, color=MUTE, align=PP_ALIGN.CENTER)


def new(title=None, kicker=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        header(s, title, kicker)
    return s


# ─────────────────────────────────────────────────────────────────────────────
# 1. Title
# ─────────────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.4))
_set(tb.text_frame.paragraphs[0],
     "FPS-Voronoi: From Mock 2-D to Real LiDAR", 36, bold=True, color=INK)
sb = s.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.9))
_set(sb.text_frame.paragraphs[0],
     "Adaptive point-cloud sampling — and measuring how much consecutive frames differ",
     20, color=MUTE)
_rule(s, Inches(0.95), Inches(4.5), Inches(6.0))
mb = s.shapes.add_textbox(Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.6))
_set(mb.text_frame.paragraphs[0],
     "Discussion deck  ·  full detail in README.md", 14, color=MUTE)

# ─────────────────────────────────────────────────────────────────────────────
# 2. Recap of the pipeline
# ─────────────────────────────────────────────────────────────────────────────
s = new("The pipeline in one slide", kicker="Context")
bullets(s, [
    (0, "Goal: pick a small sample set S that represents a large point cloud P well."),
    (0, "Each cloud point belongs to its nearest sample → implicit Voronoi cells."),
    (0, "Phase 1 — five geometric primitives (membership, covering radius, occupancy, "
        "nearest-pair, Delaunay neighbours)."),
    (0, "Phase 2 — one fused pass + three detectors:"),
    (1, "Coverage gap (cell too large) · Separation (samples too close) · Vanishing (empty cell)."),
    (0, "Phase 3 — correction loop: insert / evict with cheap one-hop updates, "
        "under an edit budget."),
    (0, "Iterated, it turns raw FPS into an adaptive resampler.", "mute"),
])

# ─────────────────────────────────────────────────────────────────────────────
# 3. Today's focus
# ─────────────────────────────────────────────────────────────────────────────
s = new("What we did today", kicker="Agenda")
bullets(s, [
    (0, "1.  Moved the whole pipeline from mock 2-D data onto real 3-D LiDAR "
        "(KITTI + nuScenes)."),
    (0, "2.  Built a temporal loop: carry the sampling across consecutive frames "
        "and measure how much it must adapt."),
    (0, "3.  Added a Chamfer quality metric — and used it to ask whether reusing "
        "the previous frame's samples is as good as rebuilding."),
    (0, ""),
    (0, "Theme for discussion: a cheap, reused sampling vs. a from-scratch one — "
        "what we gain, what we pay, and why.", "mute"),
], size=19, gap=12)

# ─────────────────────────────────────────────────────────────────────────────
# 4. Real data — key insight + loader
# ─────────────────────────────────────────────────────────────────────────────
s = new("Real LiDAR: the engine was already dimension-agnostic", kicker="Today · 1")
bullets(s, [
    (0, "Every primitive uses distance only (cdist / norm / Delaunay) → works in 2-D or 3-D unchanged."),
    (0, "So real data is a loading concern, not an algorithm change. One shared loader (data_io.py):"),
    (1, "Auto-detects format: KITTI .bin (x,y,z,intensity) · nuScenes .pcd.bin (x,y,z,intensity,ring)."),
    (1, "dims = 3 (full) or 2 (top-down, keeps the original visualizations)."),
    (1, "Optional max_range crop and min_z ground removal."),
    (0, "Why crop / remove ground: on raw LiDAR, FPS grabs sparse far/high outliers and leaves one "
        "giant central cell (~22k of 35k points).", "mute"),
], top=Inches(1.75), width=Inches(12.2), size=17, gap=7)

# ─────────────────────────────────────────────────────────────────────────────
# 5. Phase 1 on real data (image)
# ─────────────────────────────────────────────────────────────────────────────
s = new("Primitives on a real nuScenes frame (3-D, cropped)", kicker="Today · 1")
fit_image(s, os.path.join(HERE, "phase1_real.png"),
          max_w=Inches(11.6), max_h=Inches(5.0), top=Inches(1.75))
caption(s, "FPS samples · Voronoi membership · Delaunay graph · covering radii · "
           "occupancy & radius histograms — all computed on ~20.8k real points.",
        Inches(6.85))

# ─────────────────────────────────────────────────────────────────────────────
# 6. Phase 2 detectors on real data (image)
# ─────────────────────────────────────────────────────────────────────────────
s = new("The three detectors firing on real data", kicker="Today · 1")
fit_image(s, os.path.join(HERE, "phase2_real.png"),
          max_w=Inches(9.2), max_h=Inches(5.0), top=Inches(1.75))
caption(s, "Coverage gaps (insertion points) · separation violators · vanishing cells "
           "— same scene, real LiDAR.", Inches(6.85))

# ─────────────────────────────────────────────────────────────────────────────
# 7. The temporal question + method
# ─────────────────────────────────────────────────────────────────────────────
s = new("How much do consecutive frames differ?", kicker="Today · 2")
bullets(s, [
    (0, "Idea: measure the difference by how much the sampling must ADAPT, frame to frame."),
    (0, "Method (phase_3/temporal_demo.py):"),
    (1, "Frame 0 — run FPS to get samples S."),
    (1, "Each later frame — carry S forward, run Phase 3 correct() against the new cloud."),
    (1, "It patches S (insert at gaps, evict vanished/redundant) via cheap one-hop updates."),
    (0, "Headline signal: edits requested = insertions + evictions."),
    (1, "Static scene → ~0 edits · fast-changing scene → many edits."),
    (0, "10 Hz data: once the sampling settles, edits fall to 0–2 per frame.", "mute"),
], top=Inches(1.75), size=17, gap=7)

# ─────────────────────────────────────────────────────────────────────────────
# 8. Two thresholds (theory)
# ─────────────────────────────────────────────────────────────────────────────
s = new("Two threshold layers — only one decides patch vs. rebuild", kicker="Theory")
bullets(s, [
    (0, "Detector thresholds — decide HOW MANY edits a frame wants. Data-adaptive:"),
    (1, "Coverage gap  =  2.0 × the frame's own median covering radius."),
    (1, "Separation     =  0.5 × median sample-to-sample nearest distance."),
    (1, "Vanishing       =  occupancy < 1.   (each can be pinned to an absolute value.)"),
    (0, "Budget — decides PATCH vs. FULL REBUILD:"),
    (1, "If (insertions + evictions) > budget → discard S, re-run FPS from scratch."),
    (1, "A fixed integer you choose (default 40) — a compute ceiling, not derived from data."),
    (0, "Discussion point: the detector bar is relative; the budget is absolute.", "mute"),
], top=Inches(1.75), size=17, gap=7)

# ─────────────────────────────────────────────────────────────────────────────
# 9. Chamfer metric (theory)
# ─────────────────────────────────────────────────────────────────────────────
s = new("Measuring quality: Chamfer, split into two halves", kicker="Theory")
bullets(s, [
    (0, "Chamfer(P, S) = symmetric mean of nearest-neighbour distances (same as extract/)."),
    (0, "Reported as its two directed halves — they catch different failures:"),
    (1, "cover  (cloud → S):  typical cloud point's distance to nearest sample → UNDER-COVERAGE."),
    (1, "faith   (S → cloud):  typical sample's distance to nearest real point → STALE samples "
        "sitting in empty space."),
    (0, "chamfer = cover + faith   (metres)."),
    (0, "Fresh FPS has faith ≈ 0 — its samples ARE cloud points."),
    (0, "Baseline check (--baseline): rebuild fresh FPS at the SAME M; "
        "ratio = chamfer / fps_chamfer.   ratio > 1 ⇒ reuse worse than rebuild.", "mute"),
], top=Inches(1.75), size=17, gap=7)

# ─────────────────────────────────────────────────────────────────────────────
# 10. Result — plot + findings
# ─────────────────────────────────────────────────────────────────────────────
s = new("Result: reuse is sound, but ~2× the chamfer of a rebuild", kicker="Today · 3")
fit_image(s, os.path.join(HERE, "temporal_chamfer.png"),
          max_w=Inches(6.6), max_h=Inches(5.1), top=Inches(1.8),
          center_x=False, left=Inches(0.6))
bullets(s, [
    (0, "Edits settle to 0–2/frame; budget never tripped (continuous 10 Hz drive)."),
    (0, "One-hop reuse is exact — matches a full recompute; converges, doesn't diverge."),
    (0, "But chamfer ratio settles ≈ 2.0:"),
    (1, "reused S sits ~2× farther from the cloud than a fresh rebuild of equal size."),
    (0, "Why: the faith (staleness) half grows 0 → ~2 m and dominates the gap."),
], left=Inches(7.4), top=Inches(1.9), width=Inches(5.5), size=15, gap=9)

# ─────────────────────────────────────────────────────────────────────────────
# 11. Terminal output (subset table)
# ─────────────────────────────────────────────────────────────────────────────
s = new("Terminal output (KITTI, excerpt)", kicker="Today · 3")
rows = [
    "frame    pts    M  edits ins evt  cover  faith chamfer fps_cham ratio  note",
    "----------------------------------------------------------------------------",
    "    0  33863   64    -    -   -   2.514  0.001   2.515      -      -   FPS baseline",
    "    1  32720   64    0    0   0   2.566  0.745   3.311    2.530   1.31",
    "    5  36454   63    1    1   0   2.915  1.659   4.575    2.503   1.83",
    "   10  37721   54    1    0   1   3.487  2.024   5.511    2.521   2.19",
    "   15  39933   48    0    0   0   3.336  1.855   5.192    2.974   1.75",
    "   20  40925   47    1    0   1   3.268  2.171   5.439    2.416   2.25",
    "   25  40031   45    1    1   0   3.429  2.334   5.764    3.063   1.88",
    "   29  49254   39    0    0   0   3.486  1.658   5.144    2.850   1.80",
]
panel = s.shapes.add_shape(__import__("pptx").enum.shapes.MSO_SHAPE.RECTANGLE,
                           Inches(0.7), Inches(1.85), Inches(12.0), Inches(3.7))
panel.fill.solid(); panel.fill.fore_color.rgb = PANEL
panel.line.fill.background(); panel.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.95), Inches(2.0), Inches(11.6), Inches(3.4))
tf = tb.text_frame; tf.word_wrap = False
for i, line in enumerate(rows):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    _set(p, line, 12.5, font=MONO, color=INK)
    p.space_after = Pt(2)
caption(s, "edits collapse to 0–2 once settled · faith (staleness) climbs · "
           "ratio holds ≈ 2.0 throughout.", Inches(5.75))

# ─────────────────────────────────────────────────────────────────────────────
# 12. Caveat + next steps
# ─────────────────────────────────────────────────────────────────────────────
s = new("Caveat & open questions", kicker="Discussion")
bullets(s, [
    (0, "Relative thresholds can hide coarsening:"),
    (1, "as evictions outpace insertions, M drifts down (64 → ~40), cells grow,"),
    (1, "the 'gap' bar rises with them → system reports '≈0 edits' while quietly getting coarser."),
    (0, "The Chamfer metric is what exposed this — edits alone would not have."),
    (0, ""),
    (0, "Candidate fixes (should pull chamfer toward the fresh-FPS baseline):"),
    (1, "Absolute coverage threshold instead of a per-frame relative one."),
    (1, "Constant-M mode: pair each eviction with an insertion / top up with FPS."),
    (0, "Open: is the ~2× staleness cost acceptable for the compute saved? Where is the crossover?", "mute"),
], top=Inches(1.75), size=17, gap=7)

out = os.path.join(HERE, "fps_voronoi_today.pptx")
prs.save(out)
print("Saved", out, "·", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
